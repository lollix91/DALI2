#!/usr/bin/env python3
"""
DALI2 Presentation Generator
Converts Slides.xml into Slides.pptx using python-pptx with a modern template.

Usage:
    pip install python-pptx
    python Converter.py
"""

import xml.etree.ElementTree as ET
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os
import sys

# ── Color Palette (Modern Navy + Blue accent) ─────────────────
PRIMARY_DARK = RGBColor(0x0F, 0x17, 0x2A)   # Very dark navy
PRIMARY      = RGBColor(0x1E, 0x3A, 0x5F)   # Deep blue
ACCENT       = RGBColor(0x38, 0xBD, 0xF8)   # Sky blue accent
ACCENT2      = RGBColor(0x22, 0xD3, 0xEE)   # Cyan
TEXT_DARK    = RGBColor(0x1E, 0x29, 0x3B)   # Dark slate
TEXT_MEDIUM  = RGBColor(0x47, 0x55, 0x69)   # Medium gray
TEXT_LIGHT   = RGBColor(0x94, 0xA3, 0xB8)   # Light gray
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
CODE_BG      = RGBColor(0xF1, 0xF5, 0xF9)   # Slate-100
LIGHT_BG     = RGBColor(0xF8, 0xFA, 0xFC)   # Slate-50
TABLE_ALT    = RGBColor(0xEF, 0xF6, 0xFF)   # Light blue alt row
KEYWORD_CLR  = RGBColor(0x7C, 0x3A, 0xED)   # Purple for keywords
COMMENT_CLR  = RGBColor(0x64, 0x74, 0x8B)   # Gray for comments

SLIDE_WIDTH  = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

FONT_TITLE = "Segoe UI"
FONT_BODY  = "Segoe UI"
FONT_CODE  = "Consolas"


class SlideBuilder:
    """Builds a PowerPoint presentation with a modern, clean template."""

    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = SLIDE_WIDTH
        self.prs.slide_height = SLIDE_HEIGHT

    # ── Low-level helpers ──────────────────────────────────────

    def _blank(self):
        return self.prs.slides.add_slide(self.prs.slide_layouts[6])

    def _rect(self, slide, l, t, w, h, color=None):
        s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
        s.line.fill.background()
        if color:
            s.fill.solid()
            s.fill.fore_color.rgb = color
        else:
            s.fill.background()
        return s

    def _textbox(self, slide, l, t, w, h, text="",
                 size=18, color=TEXT_DARK, bold=False,
                 font=FONT_BODY, align=PP_ALIGN.LEFT):
        tb = slide.shapes.add_textbox(l, t, w, h)
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font
        p.alignment = align
        return tb

    def _bullets(self, slide, l, t, w, h, items, size=18, color=TEXT_DARK):
        tb = slide.shapes.add_textbox(l, t, w, h)
        tf = tb.text_frame
        tf.word_wrap = True
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            # Handle **bold** prefix
            if item.startswith("**") and "**" in item[2:]:
                end = item.index("**", 2)
                bold_part = item[2:end]
                rest = item[end + 2:]
                run = p.add_run()
                run.text = "\u2022 " + bold_part
                run.font.bold = True
                run.font.size = Pt(size)
                run.font.color.rgb = color
                run.font.name = FONT_BODY
                if rest:
                    r2 = p.add_run()
                    r2.text = rest
                    r2.font.size = Pt(size)
                    r2.font.color.rgb = color
                    r2.font.name = FONT_BODY
            else:
                p.text = "\u2022 " + item
                p.font.size = Pt(size)
                p.font.color.rgb = color
                p.font.name = FONT_BODY
            p.space_after = Pt(6)
        return tb

    def _code(self, slide, l, t, w, h, code, size=12):
        self._rect(slide, l, t, w, h, CODE_BG)
        # Thin left accent bar
        self._rect(slide, l, t, Inches(0.05), h, ACCENT)
        tb = slide.shapes.add_textbox(
            l + Inches(0.25), t + Inches(0.15),
            w - Inches(0.45), h - Inches(0.3)
        )
        tf = tb.text_frame
        tf.word_wrap = True
        for i, line in enumerate(code.strip().split("\n")):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = line
            p.font.size = Pt(size)
            p.font.name = FONT_CODE
            p.space_after = Pt(1)
            # Basic syntax coloring
            stripped = line.strip()
            if stripped.startswith("%%") or stripped.startswith("%"):
                p.font.color.rgb = COMMENT_CLR
            elif stripped.startswith(":-"):
                p.font.color.rgb = KEYWORD_CLR
            else:
                p.font.color.rgb = TEXT_DARK
        return tb

    def _footer(self, slide):
        self._textbox(slide,
                      Inches(0.5), SLIDE_HEIGHT - Inches(0.45),
                      Inches(6), Inches(0.3),
                      "DALI2 \u2014 University of L\u2019Aquila",
                      size=9, color=TEXT_LIGHT)

    def _slide_num(self, slide, n):
        self._textbox(slide,
                      SLIDE_WIDTH - Inches(1), SLIDE_HEIGHT - Inches(0.45),
                      Inches(0.7), Inches(0.3),
                      str(n), size=9, color=TEXT_LIGHT, align=PP_ALIGN.RIGHT)

    def _top_bar(self, slide):
        self._rect(slide, Inches(0), Inches(0),
                   SLIDE_WIDTH, Inches(0.06), PRIMARY)

    def _accent_line(self, slide, l, t, w):
        self._rect(slide, l, t, w, Inches(0.035), ACCENT)

    # ── Slide Types ────────────────────────────────────────────

    def add_title_slide(self, title, subtitle="", n=1):
        s = self._blank()
        self._rect(s, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT, PRIMARY_DARK)
        # Decorative accent lines
        self._rect(s, Inches(1), Inches(3.35), Inches(4), Inches(0.04), ACCENT)
        self._rect(s, Inches(1), Inches(3.5), Inches(2.5), Inches(0.04), ACCENT2)
        self._textbox(s, Inches(1), Inches(1.5), Inches(11), Inches(1.8),
                      title, size=48, color=WHITE, bold=True, font=FONT_TITLE)
        if subtitle:
            self._textbox(s, Inches(1), Inches(3.9), Inches(11), Inches(2),
                          subtitle, size=22,
                          color=RGBColor(0xBD, 0xD5, 0xF5))

    def add_section_slide(self, title, subtitle="", n=1):
        s = self._blank()
        self._rect(s, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT, PRIMARY_DARK)
        self._rect(s, Inches(1), Inches(4.0), Inches(3), Inches(0.04), ACCENT)
        self._textbox(s, Inches(1), Inches(2.4), Inches(11), Inches(1.5),
                      title, size=42, color=WHITE, bold=True, font=FONT_TITLE)
        if subtitle:
            self._textbox(s, Inches(1), Inches(4.4), Inches(11), Inches(1),
                          subtitle, size=20,
                          color=RGBColor(0xBD, 0xD5, 0xF5))
        self._slide_num(s, n)

    def add_content_slide(self, title, bullets, n=1, text=None):
        s = self._blank()
        self._top_bar(s)
        self._textbox(s, Inches(0.8), Inches(0.35), Inches(11.5), Inches(0.8),
                      title, size=30, color=PRIMARY_DARK, bold=True, font=FONT_TITLE)
        self._accent_line(s, Inches(0.8), Inches(1.15), Inches(2))
        y = Inches(1.5)
        if text:
            self._textbox(s, Inches(0.8), y, Inches(11.5), Inches(0.6),
                          text, size=16, color=TEXT_MEDIUM)
            y = Inches(2.2)
        if bullets:
            self._bullets(s, Inches(0.8), y, Inches(11.5), Inches(4.8),
                          bullets, size=18)
        self._footer(s)
        self._slide_num(s, n)

    def add_code_slide(self, title, code, n=1, text=None, code_size=12):
        s = self._blank()
        self._top_bar(s)
        self._textbox(s, Inches(0.8), Inches(0.35), Inches(11.5), Inches(0.8),
                      title, size=30, color=PRIMARY_DARK, bold=True, font=FONT_TITLE)
        self._accent_line(s, Inches(0.8), Inches(1.15), Inches(2))
        y = Inches(1.5)
        if text:
            self._textbox(s, Inches(0.8), y, Inches(11.5), Inches(0.5),
                          text, size=16, color=TEXT_MEDIUM)
            y = Inches(2.15)
        nlines = len(code.strip().split("\n"))
        ch = min(Inches(5.0), Inches(0.24 * nlines + 0.5))
        self._code(s, Inches(0.6), y, Inches(12), ch, code, code_size)
        self._footer(s)
        self._slide_num(s, n)

    def add_two_column_slide(self, title, left_bullets, right_content,
                             right_type="bullets", n=1, text=None):
        s = self._blank()
        self._top_bar(s)
        self._textbox(s, Inches(0.8), Inches(0.35), Inches(11.5), Inches(0.8),
                      title, size=30, color=PRIMARY_DARK, bold=True, font=FONT_TITLE)
        self._accent_line(s, Inches(0.8), Inches(1.15), Inches(2))
        y = Inches(1.5)
        if text:
            self._textbox(s, Inches(0.8), y, Inches(11.5), Inches(0.5),
                          text, size=16, color=TEXT_MEDIUM)
            y = Inches(2.15)
        if left_bullets:
            self._bullets(s, Inches(0.8), y, Inches(5.5), Inches(4.8),
                          left_bullets, size=16)
        if right_type == "code" and right_content:
            nlines = len(right_content.strip().split("\n"))
            ch = min(Inches(4.8), Inches(0.24 * nlines + 0.5))
            self._code(s, Inches(6.8), y, Inches(6), ch, right_content, 11)
        elif right_content:
            self._bullets(s, Inches(6.8), y, Inches(5.5), Inches(4.8),
                          right_content, size=16)
        self._footer(s)
        self._slide_num(s, n)

    def add_table_slide(self, title, headers, rows, n=1, text=None):
        s = self._blank()
        self._top_bar(s)
        self._textbox(s, Inches(0.8), Inches(0.35), Inches(11.5), Inches(0.8),
                      title, size=30, color=PRIMARY_DARK, bold=True, font=FONT_TITLE)
        self._accent_line(s, Inches(0.8), Inches(1.15), Inches(2))
        y = Inches(1.6)
        if text:
            self._textbox(s, Inches(0.8), y, Inches(11.5), Inches(0.5),
                          text, size=16, color=TEXT_MEDIUM)
            y = Inches(2.2)
        nr = len(rows) + 1
        nc = len(headers)
        tw = Inches(11.5)
        th = min(Inches(4.8), Inches(0.42 * nr + 0.1))
        tbl = s.shapes.add_table(nr, nc, Inches(0.8), y, tw, th).table
        for j, h in enumerate(headers):
            c = tbl.cell(0, j)
            c.text = h
            c.fill.solid()
            c.fill.fore_color.rgb = PRIMARY
            for p in c.text_frame.paragraphs:
                p.font.color.rgb = WHITE
                p.font.size = Pt(13)
                p.font.bold = True
                p.font.name = FONT_BODY
        for i, row in enumerate(rows):
            for j, val in enumerate(row):
                c = tbl.cell(i + 1, j)
                c.text = val
                c.fill.solid()
                c.fill.fore_color.rgb = WHITE if i % 2 == 0 else TABLE_ALT
                for p in c.text_frame.paragraphs:
                    p.font.size = Pt(12)
                    p.font.color.rgb = TEXT_DARK
                    p.font.name = FONT_BODY
        self._footer(s)
        self._slide_num(s, n)

    def add_diagram_slide(self, title, diagram, n=1, text=None):
        s = self._blank()
        self._top_bar(s)
        self._textbox(s, Inches(0.8), Inches(0.35), Inches(11.5), Inches(0.8),
                      title, size=30, color=PRIMARY_DARK, bold=True, font=FONT_TITLE)
        self._accent_line(s, Inches(0.8), Inches(1.15), Inches(2))
        y = Inches(1.5)
        if text:
            self._textbox(s, Inches(0.8), y, Inches(11.5), Inches(0.5),
                          text, size=16, color=TEXT_MEDIUM)
            y = Inches(2.2)
        nlines = len(diagram.strip().split("\n"))
        dh = min(Inches(4.8), Inches(0.24 * nlines + 0.5))
        self._code(s, Inches(1.5), y, Inches(10), dh, diagram, 14)
        self._footer(s)
        self._slide_num(s, n)

    def add_end_slide(self, title, subtitle="", n=1):
        s = self._blank()
        self._rect(s, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT, PRIMARY_DARK)
        self._rect(s, 0, Inches(4.2), SLIDE_WIDTH, Inches(0.04), ACCENT)
        self._rect(s, 0, Inches(4.35), SLIDE_WIDTH, Inches(0.04), ACCENT2)
        self._textbox(s, Inches(1), Inches(2.5), Inches(11.3), Inches(1.5),
                      title, size=48, color=WHITE, bold=True,
                      font=FONT_TITLE, align=PP_ALIGN.CENTER)
        if subtitle:
            self._textbox(s, Inches(1), Inches(4.7), Inches(11.3), Inches(1.5),
                          subtitle, size=22,
                          color=RGBColor(0xBD, 0xD5, 0xF5),
                          align=PP_ALIGN.CENTER)

    def save(self, path):
        self.prs.save(path)
        print(f"Saved: {path}")


# ── XML Parser ─────────────────────────────────────────────────

def parse_and_build(xml_path, output_path):
    tree = ET.parse(xml_path)
    root = tree.getroot()
    builder = SlideBuilder()
    n = 0

    for elem in root.findall("slide"):
        n += 1
        st = elem.get("type", "content")
        title = (elem.findtext("title") or "").strip()
        subtitle = (elem.findtext("subtitle") or "").strip()
        text = (elem.findtext("text") or "").strip() or None

        if st == "title":
            builder.add_title_slide(title, subtitle, n)

        elif st == "section":
            builder.add_section_slide(title, subtitle, n)

        elif st == "content":
            blist = [b.text.strip() for b in elem.findall(".//bullets/b")
                     if b.text]
            builder.add_content_slide(title, blist, n, text)

        elif st == "code":
            code = (elem.findtext("code") or "").strip()
            ce = elem.find("code")
            cs = int(ce.get("size", "12")) if ce is not None else 12
            builder.add_code_slide(title, code, n, text, cs)

        elif st == "two_column":
            le = elem.find("left")
            re = elem.find("right")
            lb = [b.text.strip() for b in (le.findall("b") if le is not None else [])
                  if b.text]
            rt = re.get("type", "bullets") if re is not None else "bullets"
            if rt == "code":
                rc = (re.findtext("code") or "").strip() if re is not None else ""
            else:
                rc = [b.text.strip() for b in (re.findall("b") if re is not None else [])
                      if b.text]
            builder.add_two_column_slide(title, lb, rc, rt, n, text)

        elif st == "table":
            te = elem.find("table")
            hdrs, rows = [], []
            if te is not None:
                he = te.find("header")
                if he is not None:
                    hdrs = [c.text.strip() for c in he.findall("cell") if c.text]
                for r in te.findall("row"):
                    rows.append([(c.text or "").strip() for c in r.findall("cell")])
            builder.add_table_slide(title, hdrs, rows, n, text)

        elif st == "diagram":
            diag = (elem.findtext("diagram") or "").strip()
            builder.add_diagram_slide(title, diag, n, text)

        elif st == "end":
            builder.add_end_slide(title, subtitle, n)

    builder.save(output_path)
    print(f"Generated {n} slides.")


if __name__ == "__main__":
    here = os.path.dirname(os.path.abspath(__file__))
    xml_in = os.path.join(here, "Slides.xml")
    pptx_out = os.path.join(here, "Slides.pptx")

    if not os.path.exists(xml_in):
        print(f"Error: {xml_in} not found!")
        sys.exit(1)

    parse_and_build(xml_in, pptx_out)
